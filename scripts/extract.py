#!/usr/bin/env python3
"""
extract.py — turn study material into plain text for the agent to read.

Supports PDF, PPTX (lecture slides — the dominant course format, and the one most
PDF-only tools miss), EPUB, TXT, and Markdown. Each format tries good libraries
first and falls back to a dependency-free path so it works on a bare machine.

Usage:
    python3 extract.py <file> [--out OUT]

Writes the extracted text to OUT (default: "<file>.extracted.txt") and prints a
small JSON metadata summary to stdout: {ok, file, format, method, chars, words,
estimated_tokens, out}. On total failure it prints {ok:false, error, hint}.

Why a script (not just "Claude, read the PDF"): extraction should be deterministic
and cheap, and binary formats need real parsing. Keep the dumb reliable part in
code; let the model do judgement on the resulting text.
"""
import argparse
import json
import os
import re
import subprocess
import sys
import zipfile


def estimate_tokens(text):
    return int(len(text.split()) / 0.75)


# ----- PDF ----------------------------------------------------------------- #
def pdf_pdftotext(path):
    try:
        out = subprocess.run(["pdftotext", "-layout", path, "-"],
                             capture_output=True, text=True, timeout=120)
        if out.returncode == 0 and out.stdout.strip():
            return out.stdout, "pdftotext"
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass
    return None, None


def pdf_pypdf(path):
    for mod in ("pypdf", "PyPDF2"):
        try:
            m = __import__(mod)
            reader = m.PdfReader(path)
            text = "\n".join((p.extract_text() or "") for p in reader.pages)
            if text.strip():
                return text, mod
        except Exception:
            continue
    return None, None


def pdf_pdfminer(path):
    try:
        from pdfminer.high_level import extract_text
        text = extract_text(path)
        if text and text.strip():
            return text, "pdfminer"
    except Exception:
        pass
    return None, None


def extract_pdf(path):
    for fn in (pdf_pdftotext, pdf_pypdf, pdf_pdfminer):
        text, method = fn(path)
        if text:
            return text, method
    raise RuntimeError("no PDF extractor available. Install one of: "
                       "poppler (pdftotext), pypdf, or pdfmin.six "
                       "(pip install pypdf).")


# ----- PPTX (lecture slides) ---------------------------------------------- #
def pptx_python_pptx(path):
    try:
        from pptx import Presentation
        prs = Presentation(path)
        chunks = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(r.text for r in para.runs)
                        if line.strip():
                            parts.append(line)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text for c in row.cells]
                        parts.append(" | ".join(cells))
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                note = slide.notes_slide.notes_text_frame.text
                if note.strip():
                    parts.append(f"[notes] {note}")
            chunks.append("\n".join(parts))
        text = "\n\n".join(chunks)
        if text.strip():
            return text, "python-pptx"
    except Exception:
        pass
    return None, None


def pptx_zip_fallback(path):
    """A .pptx is a zip; slide text lives in <a:t> elements. No deps needed."""
    try:
        with zipfile.ZipFile(path) as z:
            names = sorted(
                (n for n in z.namelist()
                 if re.match(r"ppt/slides/slide\d+\.xml$", n)),
                key=lambda n: int(re.search(r"(\d+)", n).group(1)),
            )
            chunks = []
            for i, n in enumerate(names, 1):
                xml = z.read(n).decode("utf-8", errors="ignore")
                texts = re.findall(r"<a:t>(.*?)</a:t>", xml, flags=re.S)
                texts = [re.sub(r"<[^>]+>", "", t).strip() for t in texts]
                texts = [t for t in texts if t]
                if texts:
                    chunks.append(f"--- Slide {i} ---\n" + "\n".join(texts))
            text = "\n\n".join(chunks)
            if text.strip():
                return text, "zip-xml-fallback"
    except Exception:
        pass
    return None, None


def extract_pptx(path):
    for fn in (pptx_python_pptx, pptx_zip_fallback):
        text, method = fn(path)
        if text:
            return text, method
    raise RuntimeError("could not read PPTX (file may be empty or image-only).")


# ----- EPUB ---------------------------------------------------------------- #
def extract_epub(path):
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
        book = epub.read_epub(path)
        parts = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), "html.parser")
                parts.append(soup.get_text(" ", strip=True))
        text = "\n\n".join(parts)
        if text.strip():
            return text, "ebooklib"
    except Exception:
        pass
    # stdlib fallback: epub is a zip of xhtml; strip tags crudely.
    try:
        with zipfile.ZipFile(path) as z:
            parts = []
            for n in z.namelist():
                if n.endswith((".xhtml", ".html", ".htm")):
                    raw = z.read(n).decode("utf-8", errors="ignore")
                    parts.append(re.sub(r"<[^>]+>", " ", raw))
            text = re.sub(r"[ \t]+", " ", "\n\n".join(parts))
            if text.strip():
                return text, "zip-xhtml-fallback"
    except Exception:
        pass
    raise RuntimeError("could not read EPUB.")


# ----- plain text ---------------------------------------------------------- #
def extract_text_file(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read(), "plaintext"


# ----- dispatch ------------------------------------------------------------ #
def detect_and_extract(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return extract_pdf(path)
    if ext == ".pptx":
        return extract_pptx(path)
    if ext == ".epub":
        return extract_epub(path)
    if ext in (".txt", ".md", ".markdown", ".rst"):
        return extract_text_file(path)
    # last resort: try reading as text
    return extract_text_file(path)


def main(argv=None):
    try:  # UTF-8 output regardless of console locale (e.g. GBK on Chinese Windows)
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass
    ap = argparse.ArgumentParser(description="Extract study material to plain text.")
    ap.add_argument("file")
    ap.add_argument("--out", default=None)
    args = ap.parse_args(argv)

    if not os.path.exists(args.file):
        print(json.dumps({"ok": False, "error": f"file not found: {args.file}"}))
        sys.exit(1)

    try:
        text, method = detect_and_extract(args.file)
    except Exception as e:
        print(json.dumps({"ok": False, "error": str(e),
                          "hint": "install pypdf / python-pptx / ebooklib as needed"}))
        sys.exit(1)

    out = args.out or (args.file + ".extracted.txt")
    with open(out, "w", encoding="utf-8") as f:
        f.write(text)
    print(json.dumps({
        "ok": True, "file": args.file,
        "format": os.path.splitext(args.file)[1].lower().lstrip("."),
        "method": method, "chars": len(text), "words": len(text.split()),
        "estimated_tokens": estimate_tokens(text), "out": out,
    }, ensure_ascii=False))


if __name__ == "__main__":
    main()
